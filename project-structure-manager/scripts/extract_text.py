#!/usr/bin/env python3
"""
文档文本提取工具

从PPT、Word文档中提取文本内容，保存为TXT格式，方便AI大模型读取。

用法:
    python extract_text.py <文件路径>                    # 提取单个文件
    python extract_text.py <目录路径> --batch            # 批量提取目录下所有文件
    python extract_text.py <文件路径> --output <输出目录>  # 指定输出目录
"""

import os
import sys
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


def extract_pptx_text(file_path):
    """从PPTX文件中提取文本"""
    if not PPTX_AVAILABLE:
        raise ImportError("需要安装 python-pptx: pip install python-pptx")
    
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n=== 幻灯片 {slide_num} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
                    text_content.append("")
        
        return "\n".join(text_content)
    except Exception as e:
        raise Exception(f"提取PPT文本失败: {str(e)}")


def extract_docx_text(file_path):
    """从DOCX文件中提取文本"""
    if not DOCX_AVAILABLE:
        raise ImportError("需要安装 python-docx: pip install python-docx")
    
    try:
        doc = Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # 提取表格内容
        for table in doc.tables:
            text_content.append("\n=== 表格 ===\n")
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_content.append(row_text)
            text_content.append("")
        
        return "\n".join(text_content)
    except Exception as e:
        raise Exception(f"提取Word文本失败: {str(e)}")


def extract_text(file_path, output_dir=None):
    """提取文件文本内容"""
    file_path = Path(file_path)
    
    if not file_path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    # 确定输出目录
    if output_dir:
        output_dir = Path(output_dir)
    else:
        # 默认输出到同级的txt目录
        output_dir = file_path.parent.parent / "txt"
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 根据文件类型提取
    ext = file_path.suffix.lower()
    
    if ext in ['.pptx', '.ppt']:
        text_content = extract_pptx_text(file_path)
    elif ext in ['.docx', '.doc']:
        text_content = extract_docx_text(file_path)
    else:
        raise ValueError(f"不支持的文件类型: {ext}。支持的类型: .pptx, .ppt, .docx, .doc")
    
    # 生成输出文件名
    output_filename = file_path.stem + ".txt"
    output_path = output_dir / output_filename
    
    # 保存文本文件
    output_path.write_text(text_content, encoding='utf-8')
    
    return output_path


def batch_extract(directory, output_dir=None):
    """批量提取目录下的所有文件"""
    directory = Path(directory)
    
    if not directory.is_dir():
        raise NotADirectoryError(f"不是目录: {directory}")
    
    # 支持的文件扩展名
    supported_exts = ['.pptx', '.ppt', '.docx', '.doc']
    
    files = []
    for ext in supported_exts:
        files.extend(directory.glob(f"*{ext}"))
    
    if not files:
        print(f"在 {directory} 中未找到支持的文件（.pptx, .ppt, .docx, .doc）")
        return []
    
    extracted = []
    errors = []
    
    for file_path in files:
        try:
            output_path = extract_text(file_path, output_dir)
            extracted.append((file_path, output_path))
            print(f"✓ 提取成功: {file_path.name} -> {output_path.name}")
        except Exception as e:
            errors.append((file_path, str(e)))
            print(f"✗ 提取失败: {file_path.name} - {str(e)}")
    
    print(f"\n完成: 成功 {len(extracted)} 个，失败 {len(errors)} 个")
    
    if errors:
        print("\n失败的文件:")
        for file_path, error in errors:
            print(f"  - {file_path.name}: {error}")
    
    return extracted


def main():
    parser = argparse.ArgumentParser(
        description='从PPT、Word文档中提取文本内容为TXT格式',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  %(prog)s materials/background/ppt/演示文稿.pptx
  %(prog)s materials/background/ppt/ --batch
  %(prog)s materials/background/docs/ --batch --output materials/background/txt/
        """
    )
    
    parser.add_argument(
        'path',
        type=str,
        help='文件路径或目录路径'
    )
    
    parser.add_argument(
        '--batch',
        action='store_true',
        help='批量处理目录下的所有文件'
    )
    
    parser.add_argument(
        '--output',
        type=str,
        help='输出目录（默认：同级txt目录）'
    )
    
    args = parser.parse_args()
    
    path = Path(args.path)
    
    # 检查依赖
    if not PPTX_AVAILABLE:
        print("警告: python-pptx 未安装，无法处理PPT文件")
        print("安装命令: pip install python-pptx")
    
    if not DOCX_AVAILABLE:
        print("警告: python-docx 未安装，无法处理Word文件")
        print("安装命令: pip install python-docx")
    
    if not PPTX_AVAILABLE and not DOCX_AVAILABLE:
        print("\n错误: 至少需要安装一个依赖库")
        sys.exit(1)
    
    try:
        if args.batch:
            # 批量处理
            if not path.is_dir():
                print(f"错误: {path} 不是目录")
                sys.exit(1)
            
            batch_extract(path, args.output)
        else:
            # 单个文件处理
            if not path.is_file():
                print(f"错误: {path} 不是文件")
                sys.exit(1)
            
            output_path = extract_text(path, args.output)
            print(f"✓ 提取成功: {path.name} -> {output_path}")
    
    except Exception as e:
        print(f"错误: {str(e)}")
        sys.exit(1)


if __name__ == '__main__':
    main()
